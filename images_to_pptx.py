import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image

# PowerPoint 16:9 slide size
SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)
TITLE_HEIGHT = Inches(1.0)

def is_image_file(filename):
    return filename.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff'))

def create_pptx_from_images(image_folder, output_file):
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    images = sorted([f for f in os.listdir(image_folder) if is_image_file(f)])

    if not images:
        print("No image files found in the folder.")
        return

    for image_file in images:
        image_path = os.path.join(image_folder, image_file)
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide

        # Title box (filename, no extension)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), SLIDE_WIDTH - Inches(1), TITLE_HEIGHT)
        title_frame = title_box.text_frame
        title_frame.text = os.path.splitext(image_file)[0]
        title_frame.paragraphs[0].font.size = Pt(24)

        # Load image and calculate placement/size
        try:
            img = Image.open(image_path)
            img_width, img_height = img.size
            img_ratio = img_width / img_height

            # Define max area for image under title
            max_width = SLIDE_WIDTH - Inches(1)
            max_height = SLIDE_HEIGHT - TITLE_HEIGHT - Inches(0.5)
            max_ratio = max_width / max_height

            if img_ratio > max_ratio:
                # Image is wider than the box
                display_width = max_width
                display_height = max_width / img_ratio
            else:
                # Image is taller than the box
                display_height = max_height
                display_width = max_height * img_ratio

            image_left = (SLIDE_WIDTH - display_width) / 2
            image_top = TITLE_HEIGHT + Inches(0.2) + ((max_height - display_height) / 2)

            slide.shapes.add_picture(image_path, image_left, image_top,
                                     width=display_width, height=display_height)
        except Exception as e:
            print(f"Error processing {image_file}: {e}")
            continue

    prs.save(output_file)
    print(f"Presentation saved to: {output_file}")

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python3 images_to_pptx.py <image_folder> <output.pptx>")
        sys.exit(1)

    image_dir = sys.argv[1]
    output_pptx = sys.argv[2]

    if not os.path.isdir(image_dir):
        print(f"Error: {image_dir} is not a directory.")
        sys.exit(1)

    create_pptx_from_images(image_dir, output_pptx)
